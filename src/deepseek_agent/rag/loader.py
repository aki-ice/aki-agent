from __future__ import annotations

import csv
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable


@dataclass
class Document:
    content: str
    source: str
    title: str = ""
    page: int = 0
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class DocumentParseResult:
    documents: list[Document]
    parser: str
    needs_ocr: bool = False
    warnings: list[str] = field(default_factory=list)


def load_text_file(path: str) -> Document:
    with open(path, encoding="utf-8", errors="replace") as fh:
        content = fh.read()
    return Document(content=content, source=os.path.abspath(path), title=os.path.basename(path))


def load_markdown_file(path: str) -> Document:
    with open(path, encoding="utf-8", errors="replace") as fh:
        raw = fh.read()
    text = re.sub(r"```[\s\S]*?```", " [code block] ", raw)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"!\[.*?\]\(.*?\)", "", text)
    text = re.sub(r"\[([^\]]+)\]\(.*?\)", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*{1,3}([^*]+)\*{1,3}", r"\1", text)
    return Document(content=text.strip(), source=os.path.abspath(path), title=os.path.basename(path))


def load_pdf_file(path: str) -> DocumentParseResult:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        return DocumentParseResult([], "pdf", warnings=[f"pypdf 未安装: {exc}"])
    reader = PdfReader(path)
    documents: list[Document] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        documents.append(Document(text, os.path.abspath(path), os.path.basename(path), index, {"format": "pdf", "page": index}))
    return DocumentParseResult(documents, "pypdf", needs_ocr=not any(doc.content for doc in documents), warnings=["PDF 没有可提取文本，建议使用 PaddleOCR"] if not any(doc.content for doc in documents) else [])


def load_docx_file(path: str) -> DocumentParseResult:
    try:
        from docx import Document as DocxDocument
    except ImportError as exc:
        return DocumentParseResult([], "docx", warnings=[f"python-docx 未安装: {exc}"])
    source = os.path.abspath(path)
    paragraphs = [p.text.strip() for p in DocxDocument(path).paragraphs if p.text.strip()]
    tables: list[str] = []
    for table in DocxDocument(path).tables:
        tables.append("\n".join(" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows))
    content = "\n".join(paragraphs + tables)
    return DocumentParseResult([Document(content, source, os.path.basename(path), 0, {"format": "docx"})], "python-docx", needs_ocr=not bool(content.strip()))


def load_pptx_file(path: str) -> DocumentParseResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        return DocumentParseResult([], "pptx", warnings=[f"python-pptx 未安装: {exc}"])
    source = os.path.abspath(path)
    documents: list[Document] = []
    for number, slide in enumerate(Presentation(path).slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        documents.append(Document("\n".join(texts), source, os.path.basename(path), number, {"format": "pptx", "slide": number}))
    return DocumentParseResult(documents, "python-pptx", needs_ocr=not any(doc.content for doc in documents))


def load_xlsx_file(path: str) -> DocumentParseResult:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        return DocumentParseResult([], "xlsx", warnings=[f"openpyxl 未安装: {exc}"])
    source = os.path.abspath(path)
    documents: list[Document] = []
    for sheet in load_workbook(path, read_only=True, data_only=True).worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                rows.append(" | ".join(values))
        documents.append(Document("\n".join(rows), source, os.path.basename(path), 0, {"format": "xlsx", "sheet": sheet.title}))
    return DocumentParseResult(documents, "openpyxl", needs_ocr=not any(doc.content for doc in documents))


def load_csv_file(path: str) -> DocumentParseResult:
    with open(path, encoding="utf-8", errors="replace", newline="") as fh:
        rows = [" | ".join(row) for row in csv.reader(fh)]
    return DocumentParseResult([Document("\n".join(rows), os.path.abspath(path), os.path.basename(path), 0, {"format": "csv"})], "csv")


_LOADERS: dict[str, Callable[[str], DocumentParseResult]] = {
    ".pdf": load_pdf_file,
    ".docx": load_docx_file,
    ".pptx": load_pptx_file,
    ".xlsx": load_xlsx_file,
    ".csv": load_csv_file,
}


def parse_document(path: str) -> DocumentParseResult:
    ext = Path(path).suffix.lower()
    if ext in (".md", ".markdown"):
        return DocumentParseResult([load_markdown_file(path)], "markdown")
    if ext in _LOADERS:
        return _LOADERS[ext](path)
    return DocumentParseResult([load_text_file(path)], "text")


def load_document(path: str) -> Document:
    result = parse_document(path)
    if not result.documents:
        raise RuntimeError("; ".join(result.warnings) or f"无法解析文档: {path}")
    merged = "\n\n".join(doc.content for doc in result.documents)
    metadata = {"parser": result.parser, "needs_ocr": result.needs_ocr, "parts": len(result.documents)}
    return Document(merged, os.path.abspath(path), os.path.basename(path), 0, metadata)


def load_directory(directory: str, recursive: bool = True) -> list[Document]:
    docs: list[Document] = []
    supported = {".txt", ".md", ".markdown", ".py", ".js", ".ts", ".json", ".yaml", ".yml", ".toml", ".pdf", ".docx", ".pptx", ".xlsx", ".csv"}
    for root, _dirs, files in os.walk(directory):
        for fname in files:
            if Path(fname).suffix.lower() in supported:
                try:
                    docs.append(load_document(os.path.join(root, fname)))
                except (OSError, RuntimeError):
                    pass
        if not recursive:
            break
    return docs
